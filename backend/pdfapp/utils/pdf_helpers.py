import os
import tempfile
import time
from pathlib import Path
from pypdf import PdfReader, PdfWriter
from pdf2image import convert_from_path
from PIL import Image
import img2pdf
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import pdfplumber
from io import BytesIO
import zipfile
import fitz  # PyMuPDF for advanced PDF processing

# Add poppler to PATH if it exists locally
poppler_path = os.path.join(os.path.dirname(__file__), '..', '..', 'poppler', 'poppler-23.01.0', 'Library', 'bin')
if os.path.exists(poppler_path):
    os.environ['PATH'] = os.environ.get('PATH', '') + os.pathsep + os.path.abspath(poppler_path)


class PDFProcessor:
    """Main PDF processing class with all operations"""
    
    @staticmethod
    def merge_pdfs(pdf_files):
        """Merge multiple PDF files into one"""
        try:
            writer = PdfWriter()
            
            for i, pdf_file in enumerate(pdf_files):
                try:
                    # Reset file pointer to beginning
                    pdf_file.seek(0)
                    reader = PdfReader(pdf_file)
                    
                    if len(reader.pages) == 0:
                        raise ValueError(f"PDF file {i+1} ({pdf_file.name}) has no pages")
                    
                    for page_num, page in enumerate(reader.pages):
                        try:
                            writer.add_page(page)
                        except Exception as e:
                            raise ValueError(f"Error adding page {page_num+1} from file {i+1} ({pdf_file.name}): {str(e)}")
                            
                except Exception as e:
                    if "not a PDF file" in str(e).lower() or "invalid PDF" in str(e).lower():
                        raise ValueError(f"File {i+1} ({pdf_file.name}) is not a valid PDF file")
                    else:
                        raise ValueError(f"Error processing file {i+1} ({pdf_file.name}): {str(e)}")
            
            output = BytesIO()
            writer.write(output)
            output.seek(0)
            return output
            
        except Exception as e:
            raise Exception(f"PDF merge failed: {str(e)}")
    
    @staticmethod
    def split_pdf(pdf_file, split_type='pages', split_value=None):
        """Split PDF by pages or page ranges"""
        reader = PdfReader(pdf_file)
        total_pages = len(reader.pages)
        
        if split_type == 'pages' and split_value:
            # Split by specific pages (comma-separated page numbers)
            try:
                page_numbers = [int(p.strip()) for p in split_value.split(',')]
                pages_to_extract = [p - 1 for p in page_numbers if 0 <= p - 1 < total_pages]
            except (ValueError, AttributeError):
                pages_to_extract = []
            
            writer = PdfWriter()
            for page_num in pages_to_extract:
                writer.add_page(reader.pages[page_num])
            
            output = BytesIO()
            writer.write(output)
            output.seek(0)
            return [output]
        
        elif split_type == 'each':
            # Split into individual pages
            outputs = []
            for i in range(total_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                
                output = BytesIO()
                writer.write(output)
                output.seek(0)
                outputs.append(output)
            
            return outputs
        
        elif split_type == 'range' and split_value:
            # Split by page ranges
            outputs = []
            ranges = split_value.split(',')
            
            for range_str in ranges:
                if '-' in range_str:
                    start, end = map(int, range_str.split('-'))
                    start = max(1, start) - 1
                    end = min(total_pages, end)
                else:
                    start = int(range_str) - 1
                    end = start + 1
                
                writer = PdfWriter()
                for i in range(start, end):
                    if i < total_pages:
                        writer.add_page(reader.pages[i])
                
                output = BytesIO()
                writer.write(output)
                output.seek(0)
                outputs.append(output)
            
            return outputs
        
        return []
    
    @staticmethod
    def compress_pdf(pdf_file, quality='medium'):
        """Compress PDF to reduce file size - supports both text and scanned PDFs"""
        # Get original file size
        pdf_file.seek(0, 2)
        original_size = pdf_file.tell()
        pdf_file.seek(0)
        
        # Try advanced compression first (for scanned PDFs and large files)
        if original_size > 1024 * 1024:  # Files larger than 1MB
            advanced_result = PDFProcessor.advanced_compress_pdf(pdf_file, quality)
            if advanced_result:
                return advanced_result
        
        # Fallback to basic compression for smaller files or if advanced fails
        return PDFProcessor.basic_compress_pdf(pdf_file, quality)
    
    @staticmethod
    def advanced_compress_pdf(pdf_file, quality='medium'):
        """Advanced compression using PyMuPDF for scanned PDFs and large files"""
        try:
            # Save uploaded file to temporary location for PyMuPDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_input:
                pdf_file.seek(0)
                temp_input.write(pdf_file.read())
                temp_input.flush()  # Ensure data is written
                temp_input_path = temp_input.name
            
            # Open with PyMuPDF
            doc = fitz.open(temp_input_path)
            
            # Get original size
            original_size = os.path.getsize(temp_input_path)
            
            # Compression settings based on quality
            if quality == 'low':
                # Aggressive compression for scanned documents
                image_quality = 20  # JPEG quality 0-100 (very low)
                dpi_reduction = 0.4  # Reduce resolution by 60%
                print(f"Using LOW quality: JPEG={image_quality}, DPI reduction={dpi_reduction}")
            elif quality == 'high':
                # Minimal compression - preserve quality
                image_quality = 90
                dpi_reduction = 0.9  # Reduce resolution by 10%
                print(f"Using HIGH quality: JPEG={image_quality}, DPI reduction={dpi_reduction}")
            else:  # medium
                # Balanced compression
                image_quality = 50
                dpi_reduction = 0.6  # Reduce resolution by 40%
                print(f"Using MEDIUM quality: JPEG={image_quality}, DPI reduction={dpi_reduction}")
            
            # Create new PDF with compressed images
            new_doc = fitz.open()
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Get page as image and compress it
                mat = fitz.Matrix(dpi_reduction, dpi_reduction)  # Scale down resolution
                pix = page.get_pixmap(matrix=mat)
                
                # Convert to PIL Image for compression
                img_data = pix.tobytes("png")
                pil_image = Image.open(BytesIO(img_data))
                
                # Convert to RGB if RGBA (for JPEG compression)
                if pil_image.mode == 'RGBA':
                    rgb_image = Image.new('RGB', pil_image.size, (255, 255, 255))
                    rgb_image.paste(pil_image, mask=pil_image.split()[-1])
                    pil_image = rgb_image
                
                # Additional processing for low quality
                if quality == 'low':
                    # Further reduce image size for aggressive compression
                    width, height = pil_image.size
                    new_size = (int(width * 0.8), int(height * 0.8))
                    pil_image = pil_image.resize(new_size, Image.Resampling.LANCZOS)
                
                # Compress image with quality-specific settings
                compressed_img = BytesIO()
                if quality == 'low':
                    # Use more aggressive compression for low quality
                    pil_image.save(compressed_img, format='JPEG', quality=image_quality, optimize=True, progressive=True)
                else:
                    pil_image.save(compressed_img, format='JPEG', quality=image_quality, optimize=True)
                compressed_img.seek(0)
                
                # Create new page from compressed image
                page_rect = page.rect  # Get original page dimensions
                new_page = new_doc.new_page(width=page_rect.width, height=page_rect.height)
                
                # Insert the compressed image into the new page
                compressed_img.seek(0)
                new_page.insert_image(page_rect, stream=compressed_img.read())
            
            # Save compressed PDF to a temporary file path
            temp_output_path = tempfile.mktemp(suffix='.pdf')
            new_doc.save(temp_output_path, deflate=True, clean=True)
            
            # Check compression effectiveness
            compressed_size = os.path.getsize(temp_output_path)
            compression_ratio = compressed_size / original_size
            
            print(f"Advanced compression ({quality}): {original_size} -> {compressed_size} bytes ({compression_ratio:.2%})")
            
            # Always return advanced compression result if it completed successfully
            # (Don't check effectiveness ratio - let user see the quality difference)
            with open(temp_output_path, 'rb') as f:
                output = BytesIO(f.read())
            
            # Cleanup
            new_doc.close()
            doc.close()
            
            # Cleanup temporary files with retry (Windows file handle issue)
            import time
            for attempt in range(3):
                try:
                    if os.path.exists(temp_input_path):
                        os.unlink(temp_input_path)
                    if os.path.exists(temp_output_path):
                        os.unlink(temp_output_path)
                    break
                except OSError:
                    if attempt < 2:  # Not the last attempt
                        time.sleep(0.1)  # Brief delay
                    # Ignore cleanup errors on last attempt
            
            output.seek(0)
            return output
            
        except Exception as e:
            print(f"Advanced compression failed: {e}")
            # Cleanup in case of error
            try:
                if 'doc' in locals(): doc.close()
                if 'new_doc' in locals(): new_doc.close()
                
                # Cleanup with retry for Windows
                import time
                for attempt in range(2):
                    try:
                        if 'temp_input_path' in locals() and os.path.exists(temp_input_path):
                            os.unlink(temp_input_path)
                        if 'temp_output_path' in locals() and os.path.exists(temp_output_path):
                            os.unlink(temp_output_path)
                        break
                    except OSError:
                        if attempt < 1:
                            time.sleep(0.1)
            except:
                pass
        
        return None  # Fall back to basic compression
    
    @staticmethod
    def basic_compress_pdf(pdf_file, quality='medium'):
        """Basic compression using pypdf for text-based PDFs"""
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        # Compression settings based on quality
        if quality == 'low':
            remove_annotations = True
            remove_form_fields = True
        elif quality == 'high':
            remove_annotations = False
            remove_form_fields = False
        else:  # medium
            remove_annotations = False
            remove_form_fields = False
        
        # Process each page
        for page in reader.pages:
            try:
                # Remove unnecessary elements for low quality
                if remove_annotations and '/Annots' in page:
                    try:
                        del page['/Annots']
                    except:
                        pass
                
                if remove_form_fields and '/AcroForm' in page:
                    try:
                        del page['/AcroForm']
                    except:
                        pass
                
                # Compress content streams
                if hasattr(page, 'compress_content_streams'):
                    page.compress_content_streams()
                
                writer.add_page(page)
                
            except Exception:
                # If page processing fails, add original page
                writer.add_page(page)
        
        # Apply writer-level compression
        try:
            writer.compress_identical_objects()
        except AttributeError:
            try:
                writer.remove_duplication()
            except AttributeError:
                pass
        
        # Write with compression
        output = BytesIO()
        writer.write(output)
        output.seek(0)
        return output
    
    @staticmethod
    def pdf_to_images(pdf_file, image_format='PNG', dpi=200):
        """Convert PDF pages to images - supports both text and scanned PDFs"""
        try:
            # First try with pdf2image (poppler-based) - best quality
            return PDFProcessor._pdf_to_images_pdf2image(pdf_file, image_format, dpi)
        except Exception as pdf2image_error:
            print(f"pdf2image failed, trying PyMuPDF fallback: {pdf2image_error}")
            try:
                # Fallback to PyMuPDF - works without external dependencies
                return PDFProcessor._pdf_to_images_pymupdf(pdf_file, image_format, dpi)
            except Exception as pymupdf_error:
                raise Exception(f"PDF to image conversion failed with both methods. "
                              f"pdf2image: {str(pdf2image_error)}. "
                              f"PyMuPDF: {str(pymupdf_error)}")
    
    @staticmethod
    def _pdf_to_images_pdf2image(pdf_file, image_format='PNG', dpi=200):
        """Convert PDF to images using pdf2image (poppler)"""
        # Save PDF to temporary file for pdf2image
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            pdf_file.seek(0)
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name
        
        try:
            # Convert PDF to images
            pages = convert_from_path(temp_pdf_path, dpi=dpi)
            
            # Create zip file with images
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for i, page in enumerate(pages):
                    img_buffer = BytesIO()
                    # Convert format for PIL compatibility
                    pil_format = 'JPEG' if image_format.upper() in ['JPG', 'JPEG'] else image_format
                    file_extension = 'jpg' if image_format.upper() in ['JPG', 'JPEG'] else image_format.lower()
                    
                    page.save(img_buffer, format=pil_format)
                    img_buffer.seek(0)
                    
                    zip_file.writestr(f'page_{i+1}.{file_extension}', img_buffer.getvalue())
            
            zip_buffer.seek(0)
            return zip_buffer
        
        finally:
            # Clean up temporary file
            try:
                os.unlink(temp_pdf_path)
            except:
                pass
    
    @staticmethod
    def _pdf_to_images_pymupdf(pdf_file, image_format='PNG', dpi=200):
        """Convert PDF to images using PyMuPDF (fallback method)"""
        # Save PDF to temporary file for PyMuPDF
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            pdf_file.seek(0)
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name
        
        try:
            # Open PDF with PyMuPDF
            doc = fitz.open(temp_pdf_path)
            
            # Create zip file with images
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    
                    # Calculate zoom factor for desired DPI (default is 72 DPI)
                    zoom = dpi / 72.0
                    mat = fitz.Matrix(zoom, zoom)
                    
                    # Render page to image
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert to desired format
                    if image_format.upper() == 'PNG':
                        img_data = pix.tobytes("png")
                        file_extension = 'png'
                    elif image_format.upper() in ['JPG', 'JPEG']:
                        img_data = pix.tobytes("jpeg")
                        file_extension = 'jpg'
                    else:
                        # Default to PNG for other formats
                        img_data = pix.tobytes("png")
                        file_extension = 'png'
                    
                    zip_file.writestr(f'page_{page_num+1}.{file_extension}', img_data)
            
            doc.close()
            zip_buffer.seek(0)
            return zip_buffer
        
        finally:
            # Clean up temporary file
            try:
                os.unlink(temp_pdf_path)
            except:
                pass
    
    @staticmethod
    def images_to_pdf(image_files, rotations=None):
        """Convert images to PDF - supports various image formats with optional rotations"""
        # If rotations are specified, use PyMuPDF (which supports rotation)
        if rotations and any(rot != 0 for rot in rotations):
            try:
                return PDFProcessor._images_to_pdf_pymupdf(image_files, rotations)
            except Exception as pymupdf_error:
                raise Exception(f"Images to PDF conversion with rotations failed: {str(pymupdf_error)}")
        
        # Otherwise, try img2pdf first (faster and preserves quality)
        try:
            return PDFProcessor._images_to_pdf_img2pdf(image_files)
        except Exception as img2pdf_error:
            print(f"img2pdf failed, trying PyMuPDF fallback: {img2pdf_error}")
            try:
                # Fallback to PyMuPDF method
                return PDFProcessor._images_to_pdf_pymupdf(image_files, rotations)
            except Exception as pymupdf_error:
                raise Exception(f"Images to PDF conversion failed with both methods. "
                              f"img2pdf: {str(img2pdf_error)}. "
                              f"PyMuPDF: {str(pymupdf_error)}")
    
    @staticmethod
    def _images_to_pdf_img2pdf(image_files):
        """Convert images to PDF using img2pdf library"""
        # Convert image files to bytes
        image_bytes = []
        for img_file in image_files:
            img_file.seek(0)
            image_data = img_file.read()
            
            # Validate image data
            if len(image_data) == 0:
                raise Exception("Empty image file detected")
            
            image_bytes.append(image_data)
        
        if not image_bytes:
            raise Exception("No valid images provided")
        
        # Create PDF from images
        pdf_bytes = img2pdf.convert(image_bytes)
        return BytesIO(pdf_bytes)
    
    @staticmethod
    def _images_to_pdf_pymupdf(image_files, rotations=None):
        """Convert images to PDF using PyMuPDF (fallback method with rotation support)"""
        # Create new PDF document
        doc = fitz.open()
        
        for i, img_file in enumerate(image_files):
            try:
                img_file.seek(0)
                image_data = img_file.read()
                
                if len(image_data) == 0:
                    print(f"Skipping empty image file {i+1}")
                    continue

                # Get rotation for this image (default to 0)
                rotation = rotations[i] if rotations and i < len(rotations) else 0
                
                # Open image with PyMuPDF
                img_doc = fitz.open("png", image_data)  # Works with various formats
                page = img_doc[0]
                
                # Get image dimensions
                img_rect = page.rect
                
                # Adjust dimensions for rotation (swap width/height for 90° and 270°)
                if rotation % 180 == 90:  # 90° or 270° rotation
                    pdf_page = doc.new_page(width=img_rect.height, height=img_rect.width)
                    page_rect = fitz.Rect(0, 0, img_rect.height, img_rect.width)
                else:
                    pdf_page = doc.new_page(width=img_rect.width, height=img_rect.height)
                    page_rect = img_rect
                
                # Insert image with rotation
                if rotation != 0:
                    # Normalize rotation to 0-360 range
                    rotation = rotation % 360
                    pdf_page.insert_image(page_rect, stream=image_data, rotate=rotation)
                else:
                    pdf_page.insert_image(page_rect, stream=image_data)
                
                img_doc.close()
                
            except Exception as e:
                print(f"Error processing image {i+1}: {e}")
                continue
        
        if len(doc) == 0:
            doc.close()
            raise Exception("No valid images could be processed")
        
        # Save PDF to bytes
        pdf_bytes = doc.tobytes()
        doc.close()
        
        return BytesIO(pdf_bytes)
    
    @staticmethod
    def extract_text(pdf_file):
        """Extract text from PDF with OCR fallback for scanned documents"""
        text_content = ""
        
        # First try with pdfplumber for regular text-based PDFs
        try:
            with pdfplumber.open(pdf_file) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_content += f"--- Page {page_num} ---\n{page_text}\n\n"
                
                # If we got substantial text content, return it
                if len(text_content.strip()) > 50:  # Threshold to determine if PDF has readable text
                    return text_content
        except Exception as e:
            print(f"Error with pdfplumber: {e}")
        
        # If no text found or very little text, try OCR with EasyOCR (more reliable than Tesseract)
        print("No readable text found, attempting OCR...")
        return PDFProcessor._extract_text_with_ocr(pdf_file)
    
    @staticmethod
    def _extract_text_with_ocr(pdf_file):
        """Extract text using OCR for scanned PDFs"""
        import tempfile
        import os
        from pdf2image import convert_from_bytes
        
        try:
            import easyocr
        except ImportError:
            raise Exception("OCR functionality requires 'easyocr' package. Install with: pip install easyocr")
        
        text_content = ""
        temp_files = []
        
        try:
            # Initialize EasyOCR reader (supports multiple languages)
            reader = easyocr.Reader(['en'])  # Add more languages as needed: ['en', 'es', 'fr']
            
            # Convert PDF to images
            pdf_file.seek(0)
            pdf_bytes = pdf_file.read()
            
            # Convert PDF pages to images
            images = convert_from_bytes(pdf_bytes, dpi=300)  # Higher DPI for better OCR
            
            for page_num, image in enumerate(images, 1):
                # Save image to temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                    temp_files.append(temp_img.name)
                    image.save(temp_img.name, 'PNG')
                
                # Perform OCR on the image
                try:
                    results = reader.readtext(temp_img.name)
                    
                    # Extract text from OCR results
                    page_text = ""
                    for (bbox, text, confidence) in results:
                        if confidence > 0.3:  # Only include text with reasonable confidence
                            page_text += text + " "
                    
                    if page_text.strip():
                        text_content += f"--- Page {page_num} (OCR) ---\n{page_text.strip()}\n\n"
                
                except Exception as ocr_error:
                    print(f"OCR error on page {page_num}: {ocr_error}")
                    text_content += f"--- Page {page_num} (OCR Failed) ---\nError processing this page\n\n"
        
        except Exception as e:
            print(f"OCR processing error: {e}")
            # Fallback to basic text extraction attempt
            try:
                pdf_file.seek(0)
                with pdfplumber.open(pdf_file) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text_content += f"--- Page {page_num} (Basic Extraction) ---\n"
                        page_text = page.extract_text() or "No text found on this page"
                        text_content += page_text + "\n\n"
            except:
                text_content = "Error: Unable to extract text from this PDF file."
        
        finally:
            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.unlink(temp_file)
                except:
                    pass
        
        return text_content if text_content.strip() else "No text could be extracted from this PDF."
    
    @staticmethod
    def add_watermark(pdf_file, watermark_text, position='center', opacity=0.3, 
                     font_size=36, color='gray', rotation=0, x_offset=0, y_offset=0):
        """Add enhanced text watermark to PDF with customizable options"""
        import fitz  # PyMuPDF for better watermark control
        import math
        
        # Read PDF with PyMuPDF for better control
        doc = fitz.open("pdf", pdf_file.read())
        
        # Color mapping
        color_map = {
            'red': (1, 0, 0),
            'blue': (0, 0, 1),
            'green': (0, 1, 0),
            'black': (0, 0, 0),
            'gray': (0.5, 0.5, 0.5),
            'white': (1, 1, 1),
            'yellow': (1, 1, 0),
            'orange': (1, 0.5, 0),
            'purple': (0.5, 0, 1)
        }
        
        text_color = color_map.get(color, (0.5, 0.5, 0.5))
        
        for page in doc:
            page_rect = page.rect
            width = page_rect.width
            height = page_rect.height
            
            # Calculate position based on alignment
            if position == 'center':
                x = width / 2 + x_offset
                y = height / 2 + y_offset
            elif position == 'top-left':
                x = 50 + x_offset
                y = height - 50 + y_offset
            elif position == 'top-center':
                x = width / 2 + x_offset
                y = height - 50 + y_offset
            elif position == 'top-right':
                x = width - 50 + x_offset
                y = height - 50 + y_offset
            elif position == 'middle-left':
                x = 50 + x_offset
                y = height / 2 + y_offset
            elif position == 'middle-right':
                x = width - 50 + x_offset
                y = height / 2 + y_offset
            elif position == 'bottom-left':
                x = 50 + x_offset
                y = 50 + y_offset
            elif position == 'bottom-center':
                x = width / 2 + x_offset
                y = 50 + y_offset
            else:  # bottom-right
                x = width - 50 + x_offset
                y = 50 + y_offset
            
            # Create text point
            point = fitz.Point(x, y)
            
            # Use direct text insertion with color-based transparency
            # Adjust color intensity based on opacity for a transparency effect
            if opacity < 1.0:
                # Blend color with white background based on opacity
                adjusted_color = tuple(
                    color_val * opacity + (1 - opacity) 
                    for color_val in text_color
                )
            else:
                adjusted_color = text_color
            
            # Insert text with adjusted color
            page.insert_text(
                point,
                watermark_text,
                fontsize=font_size,
                color=adjusted_color,
                overlay=True
            )
            
            # Add a semi-transparent overlay for additional opacity effect if needed
            if opacity < 0.7:  # Only for very transparent watermarks
                shape = page.new_shape()
                # Estimate text bounds
                text_len = len(watermark_text)
                estimated_width = text_len * font_size * 0.6
                estimated_height = font_size * 1.2
                
                overlay_rect = fitz.Rect(
                    x - estimated_width/2, y - estimated_height/2,
                    x + estimated_width/2, y + estimated_height/2
                )
                
                shape.draw_rect(overlay_rect)
                shape.finish(
                    fill=(1, 1, 1),  # White
                    fill_opacity=0.3 * (1 - opacity)  # Reduce opacity further
                )
                shape.commit()
        
        # Save to BytesIO
        pdf_bytes = doc.tobytes()
        doc.close()
        
        return BytesIO(pdf_bytes)
    
    @staticmethod
    def add_image_watermark(pdf_file, image_file, position='center', opacity=0.3, 
                           scale=1.0, x_offset=0, y_offset=0):
        """Add image watermark to PDF with customizable options"""
        import fitz  # PyMuPDF for better watermark control
        from PIL import Image
        import tempfile
        import os
        
        # Read PDF with PyMuPDF
        doc = fitz.open("pdf", pdf_file.read())
        
        # Create temporary file for the image
        temp_image_path = None
        try:
            # Save uploaded image to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                temp_image_path = temp_img.name
                
                # Convert image to PIL and save as PNG for consistency
                image = Image.open(image_file)
                
                # Apply opacity if needed
                if opacity < 1.0:
                    # Convert to RGBA if not already
                    if image.mode != 'RGBA':
                        image = image.convert('RGBA')
                    
                    # Apply opacity to all pixels
                    data = image.getdata()
                    new_data = []
                    for item in data:
                        if len(item) == 4:  # RGBA
                            new_data.append((item[0], item[1], item[2], int(item[3] * opacity)))
                        else:  # RGB
                            new_data.append((item[0], item[1], item[2], int(255 * opacity)))
                    
                    image.putdata(new_data)
                
                image.save(temp_img.name, 'PNG')
            
            # Get image dimensions
            img_rect = fitz.Rect(0, 0, image.width * scale, image.height * scale)
            
            for page in doc:
                page_rect = page.rect
                width = page_rect.width
                height = page_rect.height
                
                # Calculate position based on alignment
                if position == 'center':
                    x = (width - img_rect.width) / 2 + x_offset
                    y = (height - img_rect.height) / 2 + y_offset
                elif position == 'top-left':
                    x = 50 + x_offset
                    y = 50 + y_offset
                elif position == 'top-center':
                    x = (width - img_rect.width) / 2 + x_offset
                    y = 50 + y_offset
                elif position == 'top-right':
                    x = width - img_rect.width - 50 + x_offset
                    y = 50 + y_offset
                elif position == 'middle-left':
                    x = 50 + x_offset
                    y = (height - img_rect.height) / 2 + y_offset
                elif position == 'middle-right':
                    x = width - img_rect.width - 50 + x_offset
                    y = (height - img_rect.height) / 2 + y_offset
                elif position == 'bottom-left':
                    x = 50 + x_offset
                    y = height - img_rect.height - 50 + y_offset
                elif position == 'bottom-center':
                    x = (width - img_rect.width) / 2 + x_offset
                    y = height - img_rect.height - 50 + y_offset
                else:  # bottom-right
                    x = width - img_rect.width - 50 + x_offset
                    y = height - img_rect.height - 50 + y_offset
                
                # Create image rectangle at calculated position
                image_rect = fitz.Rect(x, y, x + img_rect.width, y + img_rect.height)
                
                # Insert image
                page.insert_image(image_rect, filename=temp_image_path, overlay=True)
        
        finally:
            # Clean up temporary file
            if temp_image_path and os.path.exists(temp_image_path):
                try:
                    os.unlink(temp_image_path)
                except:
                    pass
        
        # Save to BytesIO
        pdf_bytes = doc.tobytes()
        doc.close()
        
        return BytesIO(pdf_bytes)
    
    @staticmethod
    def rotate_pages(pdf_file, rotations):
        """Rotate specific pages of PDF using PyMuPDF for better compatibility
        rotations: dict with page_number: rotation_angle
        """
        try:
            # Try PyMuPDF first (better for scanned PDFs and complex content)
            return PDFProcessor._rotate_pages_pymupdf(pdf_file, rotations)
        except Exception as pymupdf_error:
            print(f"PyMuPDF rotation failed: {pymupdf_error}")
            try:
                # Fallback to pypdf
                return PDFProcessor._rotate_pages_pypdf(pdf_file, rotations)
            except Exception as pypdf_error:
                print(f"pypdf rotation failed: {pypdf_error}")
                raise Exception(f"Page rotation failed with both methods: PyMuPDF: {str(pymupdf_error)}, pypdf: {str(pypdf_error)}")

    @staticmethod
    def _rotate_pages_pymupdf(pdf_file, rotations):
        """Rotate pages using PyMuPDF (primary method - better for scanned PDFs)"""
        import fitz  # PyMuPDF
        
        # Ensure we're at the beginning and read content
        pdf_file.seek(0)
        pdf_content = pdf_file.read()
        
        # Verify we have content
        if not pdf_content:
            raise Exception("PDF file is empty or could not be read")
        
        # Open PDF with PyMuPDF
        pdf_doc = fitz.open(stream=pdf_content, filetype="pdf")
        
        try:
            # Rotate specified pages
            for page_num, rotation_angle in rotations.items():
                if 1 <= page_num <= pdf_doc.page_count:
                    page = pdf_doc[page_num - 1]  # PyMuPDF uses 0-based indexing
                    
                    # Apply rotation (normalize angle to 0, 90, 180, 270)
                    normalized_angle = rotation_angle % 360
                    if normalized_angle == 90:
                        page.set_rotation(90)
                    elif normalized_angle == 180:
                        page.set_rotation(180)
                    elif normalized_angle == 270:
                        page.set_rotation(270)
                    elif normalized_angle == 0:
                        page.set_rotation(0)
                    else:
                        # Round to nearest 90-degree increment
                        nearest_angle = round(normalized_angle / 90) * 90
                        page.set_rotation(nearest_angle % 360)
            
            # Save to BytesIO
            output = BytesIO()
            output.write(pdf_doc.write())
            output.seek(0)
            return output
            
        finally:
            pdf_doc.close()

    @staticmethod
    def _rotate_pages_pypdf(pdf_file, rotations):
        """Rotate pages using pypdf (fallback method)"""
        pdf_file.seek(0)
        pdf_content = pdf_file.read()
        
        # Verify we have content
        if not pdf_content:
            raise Exception("PDF file is empty or could not be read")
            
        # Create a new BytesIO from the content
        pdf_stream = BytesIO(pdf_content)
        reader = PdfReader(pdf_stream)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            page_num = i + 1
            if page_num in rotations:
                # Normalize rotation angle
                angle = rotations[page_num] % 360
                # pypdf expects angles in 90-degree increments
                if angle == 90:
                    page.rotate(90)
                elif angle == 180:
                    page.rotate(180)
                elif angle == 270:
                    page.rotate(270)
                # For 0 degrees or other angles, don't rotate
            writer.add_page(page)
        
        output = BytesIO()
        writer.write(output)
        output.seek(0)
        return output
    
    @staticmethod
    def secure_pdf(pdf_file, user_password=None, owner_password=None):
        """Add password protection to PDF with enhanced security"""
        try:
            # First try with PyMuPDF for better encryption support
            return PDFProcessor._secure_pdf_pymupdf(pdf_file, user_password, owner_password)
        except Exception as pymupdf_error:
            print(f"PyMuPDF secure failed: {pymupdf_error}")
            try:
                # Fallback to pypdf
                return PDFProcessor._secure_pdf_pypdf(pdf_file, user_password, owner_password)
            except Exception as pypdf_error:
                print(f"pypdf secure failed: {pypdf_error}")
                raise Exception(f"PDF encryption failed with both methods: PyMuPDF: {str(pymupdf_error)}, pypdf: {str(pypdf_error)}")

    @staticmethod
    def _secure_pdf_pymupdf(pdf_file, user_password=None, owner_password=None):
        """Secure PDF using PyMuPDF (primary method - better encryption)"""
        import fitz  # PyMuPDF
        
        # Ensure we're at the beginning and read content
        pdf_file.seek(0)
        pdf_content = pdf_file.read()
        
        if not pdf_content:
            raise Exception("PDF file is empty or could not be read")
        
        # Open PDF with PyMuPDF
        pdf_doc = fitz.open(stream=pdf_content, filetype="pdf")
        
        try:
            # Set passwords and permissions
            user_pwd = user_password or ""
            owner_pwd = owner_password or user_password or ""
            
            # Define permissions (restrict printing, copying, etc.)
            permissions = (
                fitz.PDF_PERM_PRINT |      # Allow printing
                fitz.PDF_PERM_COPY |       # Allow copying
                fitz.PDF_PERM_ANNOTATE |   # Allow annotations
                fitz.PDF_PERM_FORM |       # Allow form filling
                fitz.PDF_PERM_ACCESSIBILITY  # Allow accessibility
            )
            
            # Save with encryption
            output = BytesIO()
            output.write(pdf_doc.write(
                encryption=fitz.PDF_ENCRYPT_AES_256,  # Strong AES-256 encryption
                user_pw=user_pwd,
                owner_pw=owner_pwd,
                permissions=permissions
            ))
            output.seek(0)
            return output
            
        finally:
            pdf_doc.close()

    @staticmethod 
    def _secure_pdf_pypdf(pdf_file, user_password=None, owner_password=None):
        """Secure PDF using pypdf (fallback method)"""
        pdf_file.seek(0)
        pdf_content = pdf_file.read()
        
        if not pdf_content:
            raise Exception("PDF file is empty or could not be read")
            
        # Create a new BytesIO from the content  
        pdf_stream = BytesIO(pdf_content)
        reader = PdfReader(pdf_stream)
        writer = PdfWriter()
        
        # Copy all pages
        for page in reader.pages:
            writer.add_page(page)
        
        # Add password protection with stronger settings
        user_pwd = user_password or ""
        owner_pwd = owner_password or user_password or ""
        
        writer.encrypt(
            user_password=user_pwd,
            owner_password=owner_pwd,
            use_128bit=True  # Use 128-bit encryption
        )
        
        output = BytesIO()
        writer.write(output)
        output.seek(0)
        return output
    
    @staticmethod
    def unlock_pdf(pdf_file, password):
        """Remove password protection from PDF with enhanced support"""
        try:
            # First try with PyMuPDF for better encrypted PDF support
            return PDFProcessor._unlock_pdf_pymupdf(pdf_file, password)
        except Exception as pymupdf_error:
            print(f"PyMuPDF unlock failed: {pymupdf_error}")
            try:
                # Fallback to pypdf
                return PDFProcessor._unlock_pdf_pypdf(pdf_file, password)
            except Exception as pypdf_error:
                print(f"pypdf unlock failed: {pypdf_error}")
                raise Exception(f"PDF unlock failed with both methods: PyMuPDF: {str(pymupdf_error)}, pypdf: {str(pypdf_error)}")

    @staticmethod
    def _unlock_pdf_pymupdf(pdf_file, password):
        """Unlock PDF using PyMuPDF (primary method - better encrypted PDF support)"""
        import fitz  # PyMuPDF
        
        # Ensure we're at the beginning and read content
        pdf_file.seek(0)
        pdf_content = pdf_file.read()
        
        if not pdf_content:
            raise Exception("PDF file is empty or could not be read")
        
        # Open PDF with PyMuPDF
        pdf_doc = fitz.open(stream=pdf_content, filetype="pdf")
        
        try:
            # Check if PDF is encrypted
            if pdf_doc.needs_pass:
                # Try to authenticate
                if not pdf_doc.authenticate(password):
                    raise ValueError("Invalid password - could not decrypt PDF")
            
            # Create unlocked copy (without password protection)
            output = BytesIO()
            output.write(pdf_doc.write())  # Save without encryption
            output.seek(0)
            return output
            
        finally:
            pdf_doc.close()

    @staticmethod
    def _unlock_pdf_pypdf(pdf_file, password):
        """Unlock PDF using pypdf (fallback method)"""
        pdf_file.seek(0)
        pdf_content = pdf_file.read()
        
        if not pdf_content:
            raise Exception("PDF file is empty or could not be read")
            
        # Create a new BytesIO from the content
        pdf_stream = BytesIO(pdf_content)
        reader = PdfReader(pdf_stream)
        
        if reader.is_encrypted:
            decrypt_result = reader.decrypt(password)
            if not decrypt_result:
                raise ValueError("Invalid password - could not decrypt PDF")
        
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        
        output = BytesIO()
        writer.write(output)
        output.seek(0)
        return output

    @staticmethod
    def generate_preview_images(pdf_file):
        """Generate preview images for all PDF pages"""
        try:
            import base64
            import tempfile
            import os
            
            # Use PyMuPDF for high-quality preview generation
            doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
            preview_urls = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render page as image with good quality
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))  # 1.5x scale for better quality
                img_data = pix.tobytes("png")
                
                # Convert to base64 data URL for frontend display
                img_base64 = base64.b64encode(img_data).decode('utf-8')
                data_url = f"data:image/png;base64,{img_base64}"
                
                preview_urls.append(data_url)
            
            doc.close()
            return preview_urls
            
        except Exception as e:
            raise ValueError(f"Failed to generate preview images: {str(e)}")


class DocumentConverter:
    """Enhanced document converter with formatting preservation"""
    
    @staticmethod
    def is_pdf_scanned(pdf_file):
        """Detect if PDF is scanned (image-based) or has extractable text"""
        try:
            # Reset file pointer
            pdf_file.seek(0)
            
            doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
            total_chars = 0
            total_pages = len(doc)
            
            # Sample first few pages to determine if PDF is scanned
            pages_to_check = min(3, total_pages)
            
            for page_num in range(pages_to_check):
                page = doc.load_page(page_num)
                text = page.get_text().strip()
                total_chars += len(text)
            
            doc.close()
            
            # If very little text found across pages, likely scanned
            avg_chars_per_page = total_chars / pages_to_check if pages_to_check > 0 else 0
            is_scanned = avg_chars_per_page < 50  # More conservative threshold for scanned detection
            
            # Reset file pointer for next use
            pdf_file.seek(0)
            
            return is_scanned, avg_chars_per_page
            
        except Exception as e:
            print(f"Error detecting PDF type: {e}")
            pdf_file.seek(0)
            return False, 0  # Assume normal PDF if detection fails
    
    @staticmethod
    def enhance_image_for_ocr(img_array):
        """Pre-process image for better OCR results"""
        try:
            import cv2
            
            # Convert to grayscale
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            # Apply denoising
            denoised = cv2.fastNlMeansDenoising(gray)
            
            # Enhance contrast using CLAHE
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
            enhanced = clahe.apply(denoised)
            
            # Apply slight gaussian blur to smooth text
            smoothed = cv2.GaussianBlur(enhanced, (1, 1), 0)
            
            return smoothed
            
        except Exception as e:
            print(f"Image enhancement failed: {e}")
            return img_array  # Return original if enhancement fails
    
    @staticmethod
    def map_pdf_font_to_word(pdf_font_name):
        """Map PDF font names to Word-compatible font names"""
        # Common PDF to Word font mappings
        font_mapping = {
            # Times fonts
            'Times-Roman': 'Times New Roman',
            'Times-Bold': 'Times New Roman',
            'Times-Italic': 'Times New Roman', 
            'Times-BoldItalic': 'Times New Roman',
            'TimesNewRomanPSMT': 'Times New Roman',
            'TimesNewRoman': 'Times New Roman',
            
            # Arial/Helvetica fonts
            'Helvetica': 'Arial',
            'Helvetica-Bold': 'Arial',
            'Helvetica-Oblique': 'Arial',
            'Helvetica-BoldOblique': 'Arial',
            'ArialMT': 'Arial',
            'Arial-BoldMT': 'Arial',
            'Arial-ItalicMT': 'Arial',
            'Arial-BoldItalicMT': 'Arial',
            
            # Calibri fonts
            'Calibri': 'Calibri',
            'Calibri-Bold': 'Calibri',
            'Calibri-Italic': 'Calibri',
            'Calibri-BoldItalic': 'Calibri',
            
            # Other common fonts
            'Courier': 'Courier New',
            'Courier-Bold': 'Courier New',
            'Courier-Oblique': 'Courier New',
            'Courier-BoldOblique': 'Courier New',
            'Georgia': 'Georgia',
            'Georgia-Bold': 'Georgia',
            'Georgia-Italic': 'Georgia',
            'Georgia-BoldItalic': 'Georgia',
        }
        
        # Clean font name
        clean_font = pdf_font_name.replace('+', '').split(',')[0].split(' ')[0]
        
        # Try exact match first
        if clean_font in font_mapping:
            return font_mapping[clean_font]
        
        # Try partial matches
        for pdf_font, word_font in font_mapping.items():
            if pdf_font.lower() in clean_font.lower() or clean_font.lower() in pdf_font.lower():
                return word_font
        
        # Default fallbacks based on font characteristics
        font_lower = clean_font.lower()
        if 'times' in font_lower or 'roman' in font_lower:
            return 'Times New Roman'
        elif 'arial' in font_lower or 'helvetica' in font_lower or 'sans' in font_lower:
            return 'Arial'
        elif 'calibri' in font_lower:
            return 'Calibri'
        elif 'courier' in font_lower or 'mono' in font_lower:
            return 'Courier New'
        elif 'georgia' in font_lower:
            return 'Georgia'
        else:
            # Default to most common document font
            return 'Times New Roman'
    
    @staticmethod
    def docx_to_pdf(docx_file):
        """Convert DOCX to PDF with advanced formatting preservation"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.units import inch
            from reportlab.lib.colors import black, blue, red
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            import re
            
            # Read DOCX document
            doc = Document(docx_file)
            
            # Create PDF buffer
            buffer = BytesIO()
            pdf_doc = SimpleDocTemplate(buffer, pagesize=A4,
                                      rightMargin=72, leftMargin=72,
                                      topMargin=72, bottomMargin=18)
            
            # Get default styles and create custom ones
            styles = getSampleStyleSheet()
            
            # Register common fonts for better PDF output
            try:
                # Try to register common system fonts for better compatibility
                import platform
                system = platform.system()
                
                if system == "Windows":
                    font_paths = [
                        "C:/Windows/Fonts/times.ttf",
                        "C:/Windows/Fonts/arial.ttf", 
                        "C:/Windows/Fonts/calibri.ttf"
                    ]
                else:
                    font_paths = [
                        "/System/Library/Fonts/Times.ttc",
                        "/System/Library/Fonts/Arial.ttf",
                        "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf"
                    ]
                
                for font_path in font_paths:
                    try:
                        if os.path.exists(font_path):
                            font_name = os.path.basename(font_path).split('.')[0]
                            pdfmetrics.registerFont(TTFont(font_name, font_path))
                    except:
                        continue
            except:
                pass  # Continue with default fonts if registration fails
            
            # Create enhanced custom styles matching Word formatting
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontName='Times-Bold',
                fontSize=18,
                spaceAfter=12,
                spaceBefore=12,
                textColor=black,
                alignment=1,  # Center alignment
                leading=22
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontName='Times-Bold',
                fontSize=14,
                spaceAfter=8,
                spaceBefore=10,
                textColor=black,
                leading=17
            )
            
            normal_style = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontName='Times-Roman',
                fontSize=11,
                spaceAfter=6,
                textColor=black,
                leading=13,
                alignment=0  # Left alignment
            )
            
            # Create additional styles for different font characteristics
            small_style = ParagraphStyle(
                'SmallText',
                parent=styles['Normal'],
                fontName='Times-Roman',
                fontSize=9,
                spaceAfter=4,
                textColor=black,
                leading=11
            )
            
            large_style = ParagraphStyle(
                'LargeText',
                parent=styles['Normal'],
                fontName='Times-Roman',
                fontSize=13,
                spaceAfter=8,
                textColor=black,
                leading=16
            )
            
            # Build story for PDF
            story = []
            
            # Process document paragraphs with enhanced formatting preservation
            for para in doc.paragraphs:
                if not para.text.strip():
                    story.append(Spacer(1, 6))
                    continue
                
                # Analyze paragraph formatting from Word document
                para_text = para.text.strip()
                
                # Try to detect formatting from Word styles
                try:
                    # Get Word paragraph style information
                    word_style = para.style
                    word_font_size = None
                    word_font_name = None
                    is_bold = False
                    is_italic = False
                    
                    # Check runs for formatting
                    if para.runs:
                        # Use first run's formatting as representative
                        first_run = para.runs[0]
                        if first_run.font.size:
                            word_font_size = first_run.font.size.pt
                        if first_run.font.name:
                            word_font_name = first_run.font.name
                        is_bold = first_run.font.bold if first_run.font.bold is not None else False
                        is_italic = first_run.font.italic if first_run.font.italic is not None else False
                    
                    # Map Word font to ReportLab font
                    reportlab_font = 'Times-Roman'  # Default
                    if word_font_name:
                        if 'arial' in word_font_name.lower() or 'helvetica' in word_font_name.lower():
                            reportlab_font = 'Helvetica'
                        elif 'times' in word_font_name.lower():
                            reportlab_font = 'Times-Roman'
                        elif 'courier' in word_font_name.lower():
                            reportlab_font = 'Courier'
                    
                    # Apply bold/italic variations
                    if is_bold and is_italic:
                        if 'Helvetica' in reportlab_font:
                            reportlab_font = 'Helvetica-BoldOblique'
                        elif 'Times' in reportlab_font:
                            reportlab_font = 'Times-BoldItalic'
                        elif 'Courier' in reportlab_font:
                            reportlab_font = 'Courier-BoldOblique'
                    elif is_bold:
                        if 'Helvetica' in reportlab_font:
                            reportlab_font = 'Helvetica-Bold'
                        elif 'Times' in reportlab_font:
                            reportlab_font = 'Times-Bold'
                        elif 'Courier' in reportlab_font:
                            reportlab_font = 'Courier-Bold'
                    elif is_italic:
                        if 'Helvetica' in reportlab_font:
                            reportlab_font = 'Helvetica-Oblique'
                        elif 'Times' in reportlab_font:
                            reportlab_font = 'Times-Italic'
                        elif 'Courier' in reportlab_font:
                            reportlab_font = 'Courier-Oblique'
                    
                    # Determine font size
                    target_font_size = 11  # Default
                    if word_font_size:
                        target_font_size = max(8, min(24, int(word_font_size)))
                    
                    # Create custom style for this paragraph if needed
                    if word_font_size or word_font_name or is_bold or is_italic:
                        custom_style = ParagraphStyle(
                            f'Custom_{len(story)}',
                            parent=normal_style,
                            fontName=reportlab_font,
                            fontSize=target_font_size,
                            spaceAfter=6,
                            textColor=black,
                            leading=target_font_size * 1.2
                        )
                        story.append(Paragraph(para_text, custom_style))
                    else:
                        # Use default style detection
                        if len(para_text) < 100 and para_text.isupper():
                            story.append(Paragraph(para_text, title_style))
                        elif len(para_text) < 150 and (para_text.startswith(('Chapter', 'Section')) or re.match(r'^\d+\.', para_text)):
                            story.append(Paragraph(para_text, heading_style))
                        elif len(para_text) < 50:
                            story.append(Paragraph(para_text, small_style))
                        else:
                            story.append(Paragraph(para_text, normal_style))
                            
                except Exception:
                    # Fallback to simple style detection
                    if len(para_text) < 100 and para_text.isupper():
                        story.append(Paragraph(para_text, title_style))
                    elif len(para_text) < 150 and (para_text.startswith(('Chapter', 'Section')) or re.match(r'^\d+\.', para_text)):
                        story.append(Paragraph(para_text, heading_style))
                    else:
                        story.append(Paragraph(para_text, normal_style))
            
            # Process tables if any
            for table in doc.tables:
                # Add table data as formatted text
                story.append(Spacer(1, 12))
                story.append(Paragraph("Table Content:", heading_style))
                
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        story.append(Paragraph(row_text, normal_style))
                
                story.append(Spacer(1, 12))
            
            # Build PDF
            pdf_doc.build(story)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            # Fallback to simpler conversion
            return DocumentConverter._simple_docx_to_pdf(docx_file)
    
    @staticmethod
    def _simple_docx_to_pdf(docx_file):
        """Fallback simple DOCX to PDF conversion"""
        try:
            from docx import Document
            import docx2txt
            
            # Extract text
            try:
                text = docx2txt.process(docx_file)
            except:
                doc = Document(docx_file)
                text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
            # Create PDF
            buffer = BytesIO()
            c = canvas.Canvas(buffer, pagesize=letter)
            width, height = letter
            
            y = height - 50
            for line in text.split('\n'):
                if y < 50:
                    c.showPage()
                    y = height - 50
                
                c.drawString(50, y, line[:80])
                y -= 15
            
            c.save()
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            raise ValueError(f"Failed to convert DOCX to PDF: {str(e)}")
    
    @staticmethod
    def pdf_to_docx(pdf_file):
        """Convert PDF to DOCX with advanced formatting preservation, OCR support, and font matching"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.style import WD_STYLE_TYPE
            import numpy as np
            from PIL import Image
            import re
            
            # Detect if PDF is scanned
            is_scanned, avg_chars = DocumentConverter.is_pdf_scanned(pdf_file)
            print(f"PDF Analysis - Scanned: {is_scanned}, Avg chars per page: {avg_chars}")
            
            # Reset file pointer after detection
            pdf_file.seek(0)
            
            # Create new Word document
            doc = Document()
            
            # Set document margins to match typical PDF layout
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(1)
                section.bottom_margin = Inches(1)
                section.left_margin = Inches(1)
                section.right_margin = Inches(1)
            
            # Set up document styles with font preservation
            styles = doc.styles
            
            # Enhanced heading style with better font matching
            try:
                heading_style = styles.add_style('PDF Heading', WD_STYLE_TYPE.PARAGRAPH)
                heading_style.base_style = styles['Heading 2']
                heading_style.font.name = 'Arial'  # Common PDF heading font
                heading_style.font.size = Pt(14)
                heading_style.font.bold = True
                heading_style.font.color.rgb = RGBColor(0, 0, 0)
                heading_style.paragraph_format.space_before = Pt(12)
                heading_style.paragraph_format.space_after = Pt(6)
            except:
                heading_style = styles['Heading 2']
            
            # Enhanced normal style with comprehensive font preservation
            try:
                normal_style = styles.add_style('PDF Normal', WD_STYLE_TYPE.PARAGRAPH)
                normal_style.base_style = styles['Normal'] 
                normal_style.font.name = 'Times New Roman'  # Most common PDF body font
                normal_style.font.size = Pt(11)
                normal_style.font.color.rgb = RGBColor(0, 0, 0)
                normal_style.paragraph_format.line_spacing = 1.15
                normal_style.paragraph_format.space_after = Pt(6)
                normal_style.paragraph_format.first_line_indent = Inches(0)
            except:
                normal_style = styles['Normal']
            
            # Create font mapping for different sizes and styles found in PDF 
            font_style_map = {}
            common_fonts = ['Times New Roman', 'Arial', 'Calibri', 'Helvetica', 'Georgia']
            
            # Pre-create styles for common font sizes and variations
            for size in [8, 9, 10, 11, 12, 13, 14, 15, 16, 18, 20, 24]:
                for font_name in ['Times New Roman', 'Arial']:
                    for is_bold in [False, True]:
                        for is_italic in [False, True]:
                            try:
                                style_name = f'PDF_{font_name.replace(" ", "")}_{size}{"_B" if is_bold else ""}{"_I" if is_italic else ""}'
                                if style_name not in [s.name for s in styles]:
                                    font_style = styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
                                    font_style.base_style = styles['Normal']
                                    font_style.font.name = font_name
                                    font_style.font.size = Pt(size)
                                    font_style.font.bold = is_bold
                                    font_style.font.italic = is_italic
                                    font_style.font.color.rgb = RGBColor(0, 0, 0)
                                    font_style.paragraph_format.space_after = Pt(3 if size < 12 else 6)
                                    font_style_map[(font_name, size, is_bold, is_italic)] = font_style
                            except:
                                continue
                
            # Extract content from PDF using PyMuPDF with enhanced text extraction
            doc_fitz = fitz.open(stream=pdf_file.read(), filetype="pdf")
            
            total_pages = len(doc_fitz)
            
            for page_num in range(total_pages):
                page = doc_fitz.load_page(page_num)
                
                # Try advanced text extraction with formatting
                blocks = page.get_text("dict")
                page_text = ""
                structured_content = []
                
                # Parse text blocks with enhanced formatting information
                for block in blocks["blocks"]:
                    if "lines" in block:
                        for line in block["lines"]:
                            line_text = ""
                            line_fonts = []
                            line_spacing = line.get("bbox", [0, 0, 0, 0])
                            
                            for span in line["spans"]:
                                text = span["text"]
                                font_size = round(span["size"])
                                font_flags = span["flags"]  # Bold, italic, etc.
                                font_name = span.get("font", "Unknown")
                                
                                # Map PDF font names to Word-compatible fonts
                                word_font = DocumentConverter.map_pdf_font_to_word(font_name)
                                
                                # Extract text formatting
                                is_bold = bool(font_flags & 2**4)
                                is_italic = bool(font_flags & 2**1)
                                
                                line_text += text
                                line_fonts.append({
                                    'text': text,
                                    'size': font_size,
                                    'bold': is_bold,
                                    'italic': is_italic,
                                    'font_name': word_font,
                                    'original_font': font_name,
                                    'color': span.get("color", 0)  # Text color
                                })
                            
                            if line_text.strip():
                                # Calculate line height for spacing
                                line_height = line_spacing[3] - line_spacing[1] if len(line_spacing) >= 4 else 12
                                
                                structured_content.append({
                                    'text': line_text.strip(),
                                    'fonts': line_fonts,
                                    'bbox': line_spacing,
                                    'line_height': line_height,
                                    'x_pos': line_spacing[0] if len(line_spacing) >= 2 else 0
                                })
                                page_text += line_text + "\n"
                
                # Use OCR if PDF is scanned or if no/little text found
                should_use_ocr = is_scanned or not page_text.strip() or len(page_text.strip()) < 30
                
                if should_use_ocr:
                    try:
                        # High-resolution image conversion for better OCR
                        resolution_matrix = fitz.Matrix(4, 4) if is_scanned else fitz.Matrix(2, 2)
                        pix = page.get_pixmap(matrix=resolution_matrix)
                        img_data = pix.tobytes("png")
                        
                        # Use EasyOCR with better configuration
                        import easyocr
                        
                        # Initialize OCR reader
                        reader = easyocr.Reader(['en'], gpu=False, verbose=False)
                        
                        # Convert to image and enhance for OCR
                        img = Image.open(BytesIO(img_data))
                        img_array = np.array(img)
                        
                        # Enhance image for better OCR results
                        enhanced_img = DocumentConverter.enhance_image_for_ocr(img_array)
                        
                        # Extract text with bounding boxes for layout preservation
                        ocr_results = reader.readtext(enhanced_img, detail=1, paragraph=True,
                                                    width_ths=0.8, height_ths=0.8)
                        
                        # Sort OCR results by vertical position for better text flow
                        ocr_results.sort(key=lambda x: (x[0][0][1], x[0][0][0]))  # Sort by Y then X
                        
                        ocr_content = []
                        for result in ocr_results:
                            try:
                                # Handle different OCR result formats
                                if len(result) == 3:
                                    bbox, text, confidence = result
                                elif len(result) == 2:
                                    bbox, text = result
                                    confidence = 0.9  # Default high confidence if not provided
                                else:
                                    continue  # Skip malformed results
                                
                                if confidence > 0.5 and text.strip():  # Only use confident results
                                    # Determine if it's likely a heading based on position and text
                                    try:
                                        y_pos = bbox[0][1] if bbox and len(bbox) > 0 and len(bbox[0]) > 1 else 0
                                    except:
                                        y_pos = 0
                                    
                                    text_len = len(text.strip())
                                    
                                    is_heading = (
                                        text_len < 100 and 
                                        (text.isupper() or 
                                         re.match(r'^\d+\.|\bChapter\b|\bSection\b', text, re.I) or
                                         confidence > 0.8)
                                    )
                                    
                                    ocr_content.append({
                                        'text': text.strip(),
                                        'is_heading': is_heading,
                                        'confidence': confidence,
                                        'y_pos': y_pos
                                    })
                            except Exception as e:
                                print(f"Error processing OCR result: {e}")
                                continue
                        
                        structured_content = ocr_content
                        
                    except Exception as ocr_error:
                        print(f"OCR failed for page {page_num + 1}: {ocr_error}")
                        # Use simple text extraction as fallback
                        fallback_text = page.get_text()
                        if fallback_text.strip():
                            # Split into paragraphs for better formatting
                            paragraphs = fallback_text.split('\n\n')
                            structured_content = []
                            for para in paragraphs:
                                if para.strip():
                                    structured_content.append({
                                        'text': para.strip(),
                                        'fonts': [],
                                        'is_heading': len(para.strip()) < 100 and para.strip().isupper()
                                    })
                        else:
                            # Add empty page indicator
                            structured_content = [{'text': '[No extractable text found on this page]', 'fonts': [], 'is_heading': False}]
                
                # Add page content to Word document
                if page_num > 0:
                    # Add page break between pages
                    doc.add_page_break()
                
                # Add page header
                page_header = doc.add_heading(f'Page {page_num + 1}', level=1)
                page_header.style = heading_style
                
                # Process structured content with enhanced font preservation
                for content in structured_content:
                    text = content['text']
                    if not text.strip():
                        continue
                    
                    # Determine paragraph style and formatting
                    is_heading = content.get('is_heading', False)
                    fonts = content.get('fonts', [])
                    x_pos = content.get('x_pos', 0)
                    
                    # Analyze font characteristics for better style matching
                    if fonts:
                        # Get dominant font characteristics
                        avg_size = sum(f.get('size', 11) for f in fonts) / len(fonts)
                        has_bold = any(f.get('bold', False) for f in fonts)
                        has_italic = any(f.get('italic', False) for f in fonts)
                        dominant_font = max(set(f.get('font_name', 'Times New Roman') for f in fonts), 
                                          key=lambda x: sum(1 for f in fonts if f.get('font_name') == x))
                        
                        # Check if it looks like a heading from font analysis
                        if not is_heading:
                            is_heading = (
                                (avg_size > 13) or 
                                (has_bold and len(text) < 100) or
                                (avg_size > 12 and len(text) < 150) or
                                text.isupper() and len(text) < 200
                            )
                        
                        # Find or create appropriate style
                        target_size = max(8, min(24, round(avg_size)))
                        style_key = (dominant_font, target_size, has_bold, has_italic)
                        
                        if style_key in font_style_map:
                            target_style = font_style_map[style_key]
                        else:
                            # Use closest available style
                            if has_bold and target_size > 12:
                                target_style = heading_style
                            elif target_size <= 10:
                                target_style = font_style_map.get(('Times New Roman', 10, False, False), normal_style)
                            else:
                                target_style = normal_style
                    else:
                        # No font info available, use default logic
                        if is_heading or (len(text) < 100 and (text.isupper() or re.match(r'^\d+\.', text))):
                            target_style = heading_style
                        else:
                            target_style = normal_style
                    
                    # Add paragraph with appropriate style and formatting
                    if is_heading and avg_size > 14:
                        para = doc.add_heading(text, level=1)
                    elif is_heading:
                        para = doc.add_heading(text, level=2)
                    else:
                        para = doc.add_paragraph()
                        
                        # Create runs with specific formatting for mixed-font content
                        if fonts and len(fonts) > 1:
                            # Multiple font segments in this line
                            for font_info in fonts:
                                run = para.add_run(font_info['text'])
                                run.font.name = font_info.get('font_name', 'Times New Roman')
                                run.font.size = Pt(max(8, min(72, font_info.get('size', 11))))
                                run.font.bold = font_info.get('bold', False)
                                run.font.italic = font_info.get('italic', False)
                                run.font.color.rgb = RGBColor(0, 0, 0)
                        else:
                            # Single font for entire paragraph
                            run = para.add_run(text)
                            if fonts:
                                font_info = fonts[0]
                                run.font.name = font_info.get('font_name', 'Times New Roman')
                                run.font.size = Pt(max(8, min(72, font_info.get('size', 11))))
                                run.font.bold = font_info.get('bold', False)
                                run.font.italic = font_info.get('italic', False)
                                run.font.color.rgb = RGBColor(0, 0, 0)
                            else:
                                run.font.name = 'Times New Roman'
                                run.font.size = Pt(11)
                    
                    # Set paragraph alignment based on position
                    if x_pos > 200:  # Approximate right-aligned threshold
                        para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                    elif x_pos > 100:  # Approximate center-aligned threshold  
                        para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    else:
                        para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            doc_fitz.close()
            
            # Verify document has content
            total_paragraphs = len(doc.paragraphs)
            content_paragraphs = sum(1 for p in doc.paragraphs if p.text.strip())
            print(f"Document created - Total paragraphs: {total_paragraphs}, Content paragraphs: {content_paragraphs}")
            
            # If no content was added, add a fallback message
            if content_paragraphs <= total_pages:  # Only page headers
                doc.add_paragraph("No extractable text content found in this PDF document.")
                doc.add_paragraph("This may be due to:")
                doc.add_paragraph("• The PDF contains only images")
                doc.add_paragraph("• The text is embedded as images")
                doc.add_paragraph("• The PDF uses unsupported text encoding")
                print("Added fallback content message")
            
            # Save to BytesIO
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            print(f"Final document size: {len(buffer.getvalue())} bytes")
            return buffer
            
        except Exception as e:
            # Fallback to simpler conversion
            return DocumentConverter._simple_pdf_to_docx(pdf_file)
    
    @staticmethod
    def _simple_pdf_to_docx(pdf_file):
        """Fallback simple PDF to DOCX conversion with better content extraction"""
        try:
            from docx import Document
            
            print("Using simple PDF to DOCX conversion fallback")
            
            # Simple text extraction
            pdf_file.seek(0)
            doc_fitz = fitz.open(stream=pdf_file.read(), filetype="pdf")
            doc = Document()
            
            total_text_found = ""
            
            for page_num in range(len(doc_fitz)):
                page = doc_fitz.load_page(page_num)
                
                # Try multiple text extraction methods
                text_methods = [
                    page.get_text(),
                    page.get_text("text"),
                    page.get_text("blocks")
                ]
                
                page_text = ""
                for method_text in text_methods:
                    if isinstance(method_text, str) and method_text.strip():
                        page_text = method_text
                        break
                    elif isinstance(method_text, list):
                        # Handle blocks format
                        page_text = "\n".join([block[4] for block in method_text if len(block) > 4 and block[4].strip()])
                        if page_text.strip():
                            break
                
                # Add page content
                if page_text.strip():
                    doc.add_heading(f'Page {page_num + 1}', level=1)
                    
                    # Split into logical paragraphs
                    paragraphs = []
                    lines = page_text.split('\n')
                    current_para = ""
                    
                    for line in lines:
                        line = line.strip()
                        if not line:
                            if current_para:
                                paragraphs.append(current_para)
                                current_para = ""
                        else:
                            if current_para and len(current_para) > 500:  # Start new paragraph for long content
                                paragraphs.append(current_para)
                                current_para = line
                            else:
                                current_para += " " + line if current_para else line
                    
                    if current_para:
                        paragraphs.append(current_para)
                    
                    # Add paragraphs to document
                    for para in paragraphs:
                        if para.strip():
                            doc.add_paragraph(para.strip())
                    
                    total_text_found += page_text
                else:
                    doc.add_heading(f'Page {page_num + 1}', level=1)
                    doc.add_paragraph('[This page contains no extractable text or contains images only]')
            
            doc_fitz.close()
            
            print(f"Simple conversion extracted {len(total_text_found)} characters from {len(doc_fitz)} pages")
            
            # Add summary if no content found
            if not total_text_found.strip():
                doc.add_paragraph("No text content could be extracted from this PDF.")
                doc.add_paragraph("The PDF may contain only images or use unsupported text encoding.")
            
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            print(f"Simple conversion error: {e}")
            raise ValueError(f"Failed to convert PDF to DOCX: {str(e)}")
    
    @staticmethod
    def pdf_to_pptx(pdf_file):
        """Convert PDF to PowerPoint by extracting text and images"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            
            # Extract content from PDF
            doc_fitz = fitz.open(stream=pdf_file.read(), filetype="pdf")
            
            # Create new presentation
            prs = Presentation()
            
            for page_num in range(len(doc_fitz)):
                page = doc_fitz.load_page(page_num)
                
                # Add slide
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                title = slide.shapes.title
                title.text = f"Page {page_num + 1}"
                
                # Extract text
                text = page.get_text()
                
                # If no text found or very little text, try OCR
                if not text.strip() or len(text.strip()) < 50:
                    try:
                        # Convert page to image for OCR
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Higher resolution
                        img_data = pix.tobytes("png")
                        
                        # Use EasyOCR for text extraction
                        import easyocr
                        import numpy as np
                        from PIL import Image
                        
                        # Initialize OCR reader
                        reader = easyocr.Reader(['en'])
                        
                        # Convert to numpy array for EasyOCR
                        img = Image.open(BytesIO(img_data))
                        img_array = np.array(img)
                        
                        # Extract text using OCR
                        results = reader.readtext(img_array)
                        ocr_text = ' '.join([result[1] for result in results])
                        
                        if ocr_text.strip():
                            text = ocr_text
                        
                    except Exception as ocr_error:
                        print(f"OCR failed for page {page_num + 1}: {ocr_error}")
                        # Continue with original text (even if empty)
                
                if text.strip():
                    # Add text content
                    content = slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.text = text.strip()
                    
                    # Format text
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                else:
                    # Add placeholder text
                    content = slide.placeholders[1]
                    content.text = "[This page contains no extractable text or contains images only]"
            
            doc_fitz.close()
            
            # Save to BytesIO
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer
        
        except Exception as e:
            raise ValueError(f"Failed to convert PDF to PowerPoint: {str(e)}")
    
    @staticmethod
    def pdf_to_excel(pdf_file):
        """Convert PDF to Excel by extracting tabular data"""
        try:
            import pandas as pd
            from openpyxl import Workbook
            from openpyxl.utils.dataframe import dataframe_to_rows
            
            # Extract text from PDF
            doc_fitz = fitz.open(stream=pdf_file.read(), filetype="pdf")
            
            # Create workbook
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet
            
            for page_num in range(len(doc_fitz)):
                page = doc_fitz.load_page(page_num)
                text = page.get_text()
                
                # If no text found or very little text, try OCR
                if not text.strip() or len(text.strip()) < 50:
                    try:
                        # Convert page to image for OCR
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Higher resolution
                        img_data = pix.tobytes("png")
                        
                        # Use EasyOCR for text extraction
                        import easyocr
                        import numpy as np
                        from PIL import Image
                        
                        # Initialize OCR reader
                        reader = easyocr.Reader(['en'])
                        
                        # Convert to numpy array for EasyOCR
                        img = Image.open(BytesIO(img_data))
                        img_array = np.array(img)
                        
                        # Extract text using OCR
                        results = reader.readtext(img_array)
                        ocr_text = ' '.join([result[1] for result in results])
                        
                        if ocr_text.strip():
                            text = ocr_text
                        
                    except Exception as ocr_error:
                        print(f"OCR failed for page {page_num + 1}: {ocr_error}")
                        # Continue with original text (even if empty)
                
                # Create worksheet for each page
                ws = wb.create_sheet(f"Page_{page_num + 1}")
                
                if text.strip():
                    # Try to detect tabular data
                    lines = text.split('\n')
                    
                    # Simple table detection: lines with multiple spaces or tabs
                    table_data = []
                    for line in lines:
                        if line.strip():
                            # Split by multiple spaces or tabs
                            import re
                            cells = re.split(r'\s{2,}|\t+', line.strip())
                            if len(cells) > 1:
                                table_data.append(cells)
                            else:
                                table_data.append([line.strip()])
                    
                    # Add data to worksheet
                    for row_idx, row_data in enumerate(table_data, 1):
                        for col_idx, cell_data in enumerate(row_data, 1):
                            ws.cell(row=row_idx, column=col_idx, value=cell_data)
                else:
                    # Add placeholder
                    ws.cell(row=1, column=1, value="[No extractable text found on this page]")
            
            doc_fitz.close()
            
            # Save to BytesIO
            buffer = BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            raise ValueError(f"Failed to convert PDF to Excel: {str(e)}")
    
    @staticmethod
    def pptx_to_pdf(pptx_file):
        """Convert PowerPoint to PDF"""
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, landscape
            
            # Load presentation
            prs = Presentation(pptx_file)
            
            # Create PDF
            buffer = BytesIO()
            c = canvas.Canvas(buffer, pagesize=landscape(letter))
            width, height = landscape(letter)
            
            for slide_num, slide in enumerate(prs.slides):
                # Add slide title
                y = height - 50
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, y, f"Slide {slide_num + 1}")
                y -= 30
                
                # Extract text from slide
                c.setFont("Helvetica", 12)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines = shape.text.split('\n')
                        for line in lines:
                            if y < 50:
                                c.showPage()
                                y = height - 50
                            c.drawString(50, y, line[:100])  # Limit line length
                            y -= 15
                
                # Start new page for next slide
                if slide_num < len(prs.slides) - 1:
                    c.showPage()
            
            c.save()
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            raise ValueError(f"Failed to convert PowerPoint to PDF: {str(e)}")
    
    @staticmethod
    def excel_to_pdf(excel_file):
        """Convert Excel to PDF"""
        try:
            import pandas as pd
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors
            
            # Read Excel file
            excel_data = pd.read_excel(excel_file, sheet_name=None)  # Read all sheets
            
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=landscape(letter))
            elements = []
            
            for sheet_name, df in excel_data.items():
                # Add sheet title
                from reportlab.platypus import Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                
                styles = getSampleStyleSheet()
                title = Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading1'])
                elements.append(title)
                elements.append(Spacer(1, 12))
                
                if not df.empty:
                    # Convert DataFrame to table data
                    table_data = [df.columns.tolist()]  # Headers
                    for _, row in df.iterrows():
                        table_data.append([str(cell) for cell in row.tolist()])
                    
                    # Create table
                    table = Table(table_data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('FONTSIZE', (0, 1), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    elements.append(table)
                else:
                    no_data = Paragraph("No data in this sheet", styles['Normal'])
                    elements.append(no_data)
                
                elements.append(Spacer(1, 20))
            
            doc.build(elements)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            raise ValueError(f"Failed to convert Excel to PDF: {str(e)}")


class PDFOrganizer:
    """Handle PDF organization operations"""
    
    @staticmethod
    def organize_pdf(pdf_file, operation='auto', page_order=None):
        """Organize PDF pages based on content or user specifications"""
        try:
            doc_fitz = fitz.open(stream=pdf_file.read(), filetype="pdf")
            
            if operation == 'manual' and page_order:
                # Manual reordering based on user-specified page order
                return PDFOrganizer._manual_reorder(doc_fitz, page_order)
            elif operation == 'auto':
                # Auto-organize based on content analysis
                return PDFOrganizer._auto_organize(doc_fitz)
            elif operation == 'bookmark':
                # Organize based on bookmarks/outline
                return PDFOrganizer._organize_by_bookmarks(doc_fitz)
            elif operation == 'blank_remove':
                # Remove blank pages
                return PDFOrganizer._remove_blank_pages(doc_fitz)
            elif operation == 'duplicate_remove':
                # Remove duplicate pages
                return PDFOrganizer._remove_duplicate_pages(doc_fitz)
            else:
                raise ValueError(f"Unknown organization operation: {operation}")
                
        except Exception as e:
            raise ValueError(f"Failed to organize PDF: {str(e)}")
    
    @staticmethod
    def _manual_reorder(doc_fitz, page_order):
        """Manually reorder PDF pages based on user specification"""
        organized_doc = fitz.open()  # Create new document
        
        # Add pages in the specified order
        for page_index in page_order:
            if 0 <= page_index < len(doc_fitz):
                source_page = doc_fitz.load_page(page_index)
                organized_doc.new_page(width=source_page.rect.width, height=source_page.rect.height)
                organized_doc[-1].show_pdf_page(source_page.rect, doc_fitz, page_index)
        
        # Save to BytesIO
        buffer = BytesIO()
        organized_doc.save(buffer)
        organized_doc.close()
        doc_fitz.close()
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def _auto_organize(doc_fitz):
        """Auto-organize PDF based on content analysis"""
        organized_doc = fitz.open()  # Create new document
        
        # Analyze pages and group by content type
        text_pages = []
        image_pages = []
        mixed_pages = []
        
        for page_num in range(len(doc_fitz)):
            page = doc_fitz.load_page(page_num)
            text = page.get_text().strip()
            image_list = page.get_images()
            
            if text and image_list:
                mixed_pages.append(page_num)
            elif text:
                text_pages.append(page_num)
            elif image_list:
                image_pages.append(page_num)
        
        # Add pages in organized order: text first, then mixed, then images
        for page_num in text_pages + mixed_pages + image_pages:
            organized_doc.insert_pdf(doc_fitz, from_page=page_num, to_page=page_num)
        
        # Save to BytesIO
        buffer = BytesIO()
        organized_doc.save(buffer)
        organized_doc.close()
        doc_fitz.close()
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def _organize_by_bookmarks(doc_fitz):
        """Organize PDF based on bookmarks/outline"""
        toc = doc_fitz.get_toc()  # Get table of contents
        
        if not toc:
            # No bookmarks found, return original
            buffer = BytesIO()
            doc_fitz.save(buffer)
            buffer.seek(0)
            return buffer
        
        organized_doc = fitz.open()
        
        # Sort TOC entries by page number
        toc.sort(key=lambda x: x[2])  # Sort by page number (index 2)
        
        processed_pages = set()
        for level, title, page_num in toc:
            if page_num not in processed_pages:
                organized_doc.insert_pdf(doc_fitz, from_page=page_num-1, to_page=page_num-1)
                processed_pages.add(page_num)
        
        # Add any remaining pages
        for page_num in range(len(doc_fitz)):
            if page_num not in processed_pages:
                organized_doc.insert_pdf(doc_fitz, from_page=page_num, to_page=page_num)
        
        buffer = BytesIO()
        organized_doc.save(buffer)
        organized_doc.close()
        doc_fitz.close()
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def _remove_blank_pages(doc_fitz):
        """Remove blank or nearly blank pages"""
        organized_doc = fitz.open()
        
        for page_num in range(len(doc_fitz)):
            page = doc_fitz.load_page(page_num)
            text = page.get_text().strip()
            image_list = page.get_images()
            
            # Keep page if it has substantial text or images
            if len(text) > 50 or image_list:
                organized_doc.insert_pdf(doc_fitz, from_page=page_num, to_page=page_num)
        
        buffer = BytesIO()
        organized_doc.save(buffer)
        organized_doc.close()
        doc_fitz.close()
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def _remove_duplicate_pages(doc_fitz):
        """Remove duplicate pages based on text content"""
        organized_doc = fitz.open()
        seen_content = set()
        
        for page_num in range(len(doc_fitz)):
            page = doc_fitz.load_page(page_num)
            text = page.get_text().strip()
            
            # Create content hash
            content_hash = hash(text) if text else hash(str(page.get_images()))
            
            if content_hash not in seen_content:
                organized_doc.insert_pdf(doc_fitz, from_page=page_num, to_page=page_num)
                seen_content.add(content_hash)
        
        buffer = BytesIO()
        organized_doc.save(buffer)
        organized_doc.close()
        doc_fitz.close()
        buffer.seek(0)
        return buffer